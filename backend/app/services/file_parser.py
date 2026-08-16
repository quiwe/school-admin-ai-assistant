from pathlib import Path
import re
import zipfile

import olefile
import xlrd
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


class FileParseError(ValueError):
    pass


SUPPORTED_EXTENSIONS = {".pdf", ".doc", ".docx", ".ppt", ".pptx", ".txt", ".xls", ".xlsx"}

# 解析上限，防止恶意/异常文件导致内存耗尽或输出失控。
MAX_PDF_PAGES = 200
MAX_SHEET_ROWS = 5000
MAX_SHEET_COLS = 100
MAX_ZIP_UNCOMPRESSED_BYTES = 200 * 1024 * 1024  # 200 MB
MAX_LEGACY_STREAM_BYTES = 20 * 1024 * 1024  # 20 MB
MAX_TXT_BYTES = 10 * 1024 * 1024  # 10 MB

TRUNCATED_MARKER = "（内容过多，已截断）"


def _check_zip_safety(path: Path) -> None:
    """Office 新格式是 ZIP 包，解压前检查总解压体积，防止 zip 炸弹。"""
    try:
        with zipfile.ZipFile(str(path)) as archive:
            total = sum(info.file_size for info in archive.infolist())
    except zipfile.BadZipFile as exc:
        raise FileParseError(f"文件不是有效的 Office 文档：{exc}") from exc
    if total > MAX_ZIP_UNCOMPRESSED_BYTES:
        raise FileParseError("文件解压后体积过大，可能为异常压缩文件，请检查后重新上传。")


def parse_file(path: Path) -> str:
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            return parse_pdf(path)
        if suffix == ".doc":
            return parse_legacy_office(path, "DOC")
        if suffix == ".docx":
            return parse_docx(path)
        if suffix == ".ppt":
            return parse_legacy_office(path, "PPT")
        if suffix == ".pptx":
            return parse_pptx(path)
        if suffix == ".xls":
            return parse_xls(path)
        if suffix == ".xlsx":
            return parse_xlsx(path)
        if suffix == ".txt":
            data = path.read_bytes()
            if len(data) > MAX_TXT_BYTES:
                raise FileParseError("TXT 文件过大，请上传不超过 10 MB 的文本文件。")
            return data.decode("utf-8-sig", errors="replace")
    except FileParseError:
        raise
    except Exception as exc:
        raise FileParseError(f"文件解析失败：{exc}") from exc
    raise FileParseError("暂不支持该文件类型，请上传 PDF、DOC、DOCX、PPT、PPTX、TXT、XLS 或 XLSX 文件。")


def parse_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        if index > MAX_PDF_PAGES:
            pages.append(TRUNCATED_MARKER)
            break
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"第{index}页\n{text.strip()}")
    return "\n\n".join(pages).strip()


def parse_docx(path: Path) -> str:
    _check_zip_safety(path)
    doc = Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return "\n".join(paragraphs).strip()


def parse_pptx(path: Path) -> str:
    _check_zip_safety(path)
    presentation = Presentation(str(path))
    lines: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                text = "\n".join(
                    paragraph.text.strip()
                    for paragraph in shape.text_frame.paragraphs
                    if paragraph.text.strip()
                )
                if text:
                    slide_lines.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_lines.append(" | ".join(cells))
        if slide_lines:
            lines.append(f"第{slide_index}页")
            lines.extend(slide_lines)
    return "\n".join(lines).strip()


def parse_xls(path: Path) -> str:
    workbook = xlrd.open_workbook(str(path))
    lines: list[str] = []
    for sheet in workbook.sheets():
        lines.append(f"工作表：{sheet.name}")
        col_count = min(sheet.ncols, MAX_SHEET_COLS)
        row_count = min(sheet.nrows, MAX_SHEET_ROWS)
        for row_index in range(row_count):
            values = [
                format_xls_cell(workbook, sheet.cell(row_index, col_index))
                for col_index in range(col_count)
            ]
            values = [value for value in values if value]
            if values:
                lines.append(" | ".join(values))
        if sheet.nrows > MAX_SHEET_ROWS:
            lines.append(TRUNCATED_MARKER)
    return "\n".join(lines).strip()


def parse_xlsx(path: Path) -> str:
    _check_zip_safety(path)
    workbook = load_workbook(str(path), data_only=True, read_only=True)
    try:
        lines: list[str] = []
        for sheet in workbook.worksheets:
            lines.append(f"工作表：{sheet.title}")
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if row_count >= MAX_SHEET_ROWS:
                    lines.append(TRUNCATED_MARKER)
                    break
                values = [
                    str(value).strip()
                    for value in row[:MAX_SHEET_COLS]
                    if value is not None and str(value).strip()
                ]
                if values:
                    lines.append(" | ".join(values))
                row_count += 1
        return "\n".join(lines).strip()
    finally:
        workbook.close()


def parse_legacy_office(path: Path, label: str) -> str:
    if not olefile.isOleFile(str(path)):
        raise FileParseError(f"{label} 文件不是有效的旧版 Office OLE 文件。")

    chunks: list[str] = []
    with olefile.OleFileIO(str(path)) as ole:
        for entry in ole.listdir(streams=True, storages=False):
            stream_name = "/".join(entry)
            try:
                with ole.openstream(entry) as stream:
                    data = stream.read(MAX_LEGACY_STREAM_BYTES + 1)
            except Exception:
                continue
            if len(data) > MAX_LEGACY_STREAM_BYTES:
                continue
            text = extract_readable_text(data)
            if text:
                chunks.append(f"{stream_name}\n{text}")

    result = "\n\n".join(chunks).strip()
    if not result:
        raise FileParseError(
            f"旧版 {label} 文件未提取到可读文本。可尝试用 Office/WPS 另存为 DOCX、PPTX、PDF 或 TXT 后再上传。"
        )
    return result


def extract_readable_text(data: bytes) -> str:
    lines: list[str] = []
    for encoding in ("utf-16le", "gb18030", "utf-8"):
        try:
            decoded = data.decode(encoding, errors="ignore")
        except Exception:
            continue
        lines.extend(readable_lines(decoded))

    deduped: list[str] = []
    seen: set[str] = set()
    for line in lines:
        if line not in seen:
            seen.add(line)
            deduped.append(line)
    return "\n".join(deduped[:800]).strip()


def readable_lines(text: str) -> list[str]:
    normalized = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]+", "\n", text)
    normalized = re.sub(r"[ \t]{2,}", " ", normalized)
    lines: list[str] = []
    for line in normalized.splitlines():
        clean = line.strip()
        if is_readable_line(clean):
            lines.append(clean)
    return lines


def is_readable_line(value: str) -> bool:
    if len(value) < 2:
        return False
    useful = sum(1 for char in value if char.isalnum() or "\u4e00" <= char <= "\u9fff")
    if useful < 2:
        return False
    return useful / max(len(value), 1) >= 0.25


def format_xls_cell(workbook, cell) -> str:
    if cell.ctype == xlrd.XL_CELL_EMPTY:
        return ""
    if cell.ctype == xlrd.XL_CELL_DATE:
        try:
            return str(xlrd.xldate.xldate_as_datetime(cell.value, workbook.datemode)).strip()
        except Exception:
            return str(cell.value).strip()
    if cell.ctype == xlrd.XL_CELL_NUMBER:
        number = cell.value
        if float(number).is_integer():
            return str(int(number))
        return str(number).strip()
    return str(cell.value).strip()


def extract_faq_rows_from_spreadsheet(path: Path) -> list[dict[str, str]]:
    if path.suffix.lower() == ".xls":
        return extract_faq_rows_from_xls(path)
    return extract_faq_rows_from_xlsx(path)


def extract_faq_rows_from_xlsx(path: Path) -> list[dict[str, str]]:
    _check_zip_safety(path)
    workbook = load_workbook(str(path), data_only=True, read_only=True)
    rows: list[dict[str, str]] = []
    try:
        for sheet in workbook.worksheets:
            question_index = answer_index = category_index = None
            row_count = 0
            for raw_row in sheet.iter_rows(values_only=True):
                if row_count >= MAX_SHEET_ROWS:
                    break
                values = tuple(raw_row[:MAX_SHEET_COLS])
                if row_count == 0:
                    headers = [str(cell).strip() if cell is not None else "" for cell in values]
                    question_index = find_header(headers, ["问题", "question"])
                    answer_index = find_header(headers, ["答案", "answer"])
                    category_index = find_header(headers, ["分类", "category"])
                    if question_index is None or answer_index is None:
                        break
                    row_count += 1
                    continue
                question = safe_cell(values, question_index)
                answer = safe_cell(values, answer_index)
                if question and answer:
                    rows.append(
                        {
                            "question": question,
                            "answer": answer,
                            "category": safe_cell(values, category_index) if category_index is not None else "其他",
                        }
                    )
                row_count += 1
        return rows
    finally:
        workbook.close()


def extract_faq_rows_from_xls(path: Path) -> list[dict[str, str]]:
    workbook = xlrd.open_workbook(str(path))
    rows: list[dict[str, str]] = []
    for sheet in workbook.sheets():
        if sheet.nrows == 0:
            continue
        headers = [format_xls_cell(workbook, sheet.cell(0, index)) for index in range(sheet.ncols)]
        question_index = find_header(headers, ["问题", "question"])
        answer_index = find_header(headers, ["答案", "answer"])
        category_index = find_header(headers, ["分类", "category"])
        if question_index is None or answer_index is None:
            continue
        for row_index in range(1, sheet.nrows):
            raw_row = tuple(
                format_xls_cell(workbook, sheet.cell(row_index, col_index)) for col_index in range(sheet.ncols)
            )
            question = safe_cell(raw_row, question_index)
            answer = safe_cell(raw_row, answer_index)
            if question and answer:
                rows.append(
                    {
                        "question": question,
                        "answer": answer,
                        "category": safe_cell(raw_row, category_index) if category_index is not None else "其他",
                    }
                )
    return rows


def find_header(headers: list[str], candidates: list[str]) -> int | None:
    normalized = [header.lower() for header in headers]
    for candidate in candidates:
        if candidate.lower() in normalized:
            return normalized.index(candidate.lower())
    return None


def safe_cell(row: tuple, index: int) -> str:
    if index >= len(row) or row[index] is None:
        return ""
    return str(row[index]).strip()
